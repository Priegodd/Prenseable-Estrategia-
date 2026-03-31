
from __future__ import annotations

import argparse
import re
import sys
import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Iterable

from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
INPUT_DIR = BASE_DIR / "input"
OUTPUT_DIR = BASE_DIR / "output"
BRAND_DIR = BASE_DIR / "assets" / "brand-kit"

COLOR_PRIMARY = RGBColor(0xFF, 0x40, 0xB4)
COLOR_LIGHT = RGBColor(0xEC, 0xEC, 0xEC)
COLOR_MUTED = RGBColor(0x72, 0x72, 0x76)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_BRAND_DARK = RGBColor(0x5A, 0x76, 0x7E)
COLOR_BRAND_TEAL = RGBColor(0x54, 0xA7, 0xC5)
COLOR_BRAND_PALE = RGBColor(0xF7, 0xF7, 0xF8)
COLOR_CARD_LINE = RGBColor(0xCC, 0xD4, 0xD9)

TIER_COLORS = {
    "Tier 1": COLOR_PRIMARY,
    "Tier 2": COLOR_MUTED,
    "Tier 3": COLOR_LIGHT,
    "Tier 4": COLOR_DARK,
}

def resolve_asset(*relative_parts: str) -> Path:
    candidates = [
        BRAND_DIR.joinpath(*relative_parts),
        BASE_DIR.joinpath(*relative_parts),
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0]


ASSET_LOGO = resolve_asset("PRS_variables", "PRS_logo.png")
ASSET_SKYLINE = resolve_asset("Mesa de trabajo 2-8.png")
ASSET_TAGLINE = resolve_asset("Prenseables_bajada", "Mesa de trabajo 168-8.png")
ASSET_BRACKET_PINK = resolve_asset("Corchetes", "Conchetes_rosa", "Mesa de trabajo 175-8.png")
ASSET_BRACKET_WHITE = resolve_asset("Corchetes", "Conchetes_blanco", "Mesa de trabajo 176-8.png")


@dataclass
class Competitor:
    name: str
    total: str
    tiers: dict[str, int] = field(default_factory=dict)


@dataclass
class Persona:
    archetype: str = ""
    age: str = ""
    profession: str = ""
    context: str = ""
    pains: str = ""
    goals: str = ""
    personality: list[str] = field(default_factory=list)
    info_behavior: str = ""


@dataclass
class Pillar:
    name: str
    description: str


@dataclass
class Journalist:
    name: str
    outlet: str
    beat: str


@dataclass
class SocialChannel:
    platform: str
    followers: str
    frequency: str
    tone: str
    observations: str


@dataclass
class Tactic:
    name: str
    description: str


@dataclass
class ParsedData:
    raw_text: str
    general: dict[str, str] = field(default_factory=dict)
    competitors: list[Competitor] = field(default_factory=list)
    competitor_conclusions: str = ""
    concept: dict[str, str] = field(default_factory=dict)
    concept_themes: list[str] = field(default_factory=list)
    persona: Persona = field(default_factory=Persona)
    pillars: list[Pillar] = field(default_factory=list)
    monthly_topics: list[tuple[str, list[str]]] = field(default_factory=list)
    journalists: list[Journalist] = field(default_factory=list)
    strategic_steps: list[str] = field(default_factory=list)
    socials: list[SocialChannel] = field(default_factory=list)
    social_summary: str = ""
    tactics: list[Tactic] = field(default_factory=list)
    execution_plan: str = ""
    metrics: list[str] = field(default_factory=list)
    pending: list[str] = field(default_factory=list)


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix == ".txt":
            return path.read_text(encoding="utf-8")
        if suffix == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if suffix == ".pdf":
            reader = PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
    except UnicodeDecodeError as exc:
        raise ValueError("El archivo TXT no esta en UTF-8 o no se pudo leer correctamente.") from exc
    except Exception as exc:
        raise ValueError(f"No se pudo leer el archivo {path.name}. Revisa que sea un {suffix} valido.") from exc
    raise ValueError(f"Formato no soportado: {suffix}")


def normalize_label(label: str) -> str:
    label = label.strip().lower()
    replacements = {
        "á": "a",
        "é": "e",
        "í": "i",
        "ó": "o",
        "ú": "u",
        "ñ": "n",
        "/": " ",
        "-": " ",
    }
    for source, target in replacements.items():
        label = label.replace(source, target)
    return re.sub(r"\s+", " ", label).strip()


def strip_bullet(line: str) -> str:
    return re.sub(r"^[\-\*\u2022]\s*", "", line).strip()


def split_label_value(line: str) -> tuple[str, str] | None:
    clean = strip_bullet(line)
    if ":" not in clean:
        return None
    label, value = clean.split(":", 1)
    return label.strip(), value.strip()


def split_multi_value(value: str) -> list[str]:
    return [item.strip() for item in re.split(r"[;,\|]", value) if item.strip()]


def safe_int(value: str) -> int:
    digits = re.sub(r"[^\d]", "", value)
    return int(digits) if digits else 0


def parse_document(text: str) -> ParsedData:
    data = ParsedData(raw_text=text)
    section = ""

    for raw_line in text.splitlines():
        line = raw_line.strip().replace("\uf0b7", "-")
        if not line:
            continue
        if line.upper() == line and len(line) > 3 and not line.startswith("-"):
            section = normalize_label(line)
            continue

        if section == "datos generales":
            parse_general_line(line, data)
        elif section == "analisis de competidores":
            parse_competitor_line(line, data)
        elif section == "analisis de concepto clave":
            parse_concept_line(line, data)
        elif section == "buyer persona":
            parse_persona_line(line, data)
        elif section == "pilares de comunicacion":
            parse_pillar_line(line, data)
        elif section == "propuestas de tematicas":
            parse_topic_line(line, data)
        elif section == "periodistas clave":
            parse_journalist_line(line, data)
        elif section == "pasos estrategicos":
            parse_step_line(line, data)
        elif section == "estado de redes sociales del cliente":
            parse_social_line(line, data)
        elif section == "tacticas tipos de contenido":
            parse_tactic_line(line, data)
        elif section == "plan de ejecucion":
            parse_execution_line(line, data)
        elif section == "metricas de exito":
            parse_metric_line(line, data)

    post_process(data)
    return data


def parse_general_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if pair:
        label, value = pair
        data.general[normalize_label(label)] = value


def parse_competitor_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    normalized = normalize_label(label)
    if normalized.startswith("competidor"):
        parts = [part.strip() for part in value.split("|") if part.strip()]
        name = parts[0] if parts else "[COMPLETAR: nombre del competidor]"
        total = "[COMPLETAR: total de notas]"
        tiers: dict[str, int] = {}
        for part in parts[1:]:
            if ":" not in part:
                continue
            key, current = part.split(":", 1)
            key = normalize_label(key)
            if key == "total de notas":
                total = current.strip()
            elif key.startswith("tier"):
                tier_number = re.sub(r"[^\d]", "", key) or "?"
                tiers[f"Tier {tier_number}"] = safe_int(current)
        data.competitors.append(Competitor(name=name, total=total, tiers=tiers))
    elif normalized == "conclusiones del analisis":
        data.competitor_conclusions = value


def parse_concept_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    normalized = normalize_label(label)
    if normalized == "tematicas mas repetidas":
        data.concept_themes = split_multi_value(value)
    else:
        data.concept[normalized] = value


def parse_persona_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    normalized = normalize_label(label)
    persona = data.persona
    mapping = {
        "nombre o arquetipo": "archetype",
        "edad": "age",
        "profesion": "profession",
        "escenario en el que vive (contexto)": "context",
        "dolores frustraciones": "pains",
        "metas y objetivos": "goals",
        "comportamiento para buscar informacion": "info_behavior",
    }
    if normalized == "personalidad (3 conceptos)":
        persona.personality = split_multi_value(value)
    elif normalized in mapping:
        setattr(persona, mapping[normalized], value)


def parse_pillar_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    if normalize_label(label).startswith("pilar"):
        parts = [part.strip() for part in value.split("|", 1)]
        data.pillars.append(
            Pillar(
                name=parts[0] if parts else "[COMPLETAR: nombre del pilar]",
                description=parts[1] if len(parts) > 1 else "[COMPLETAR: descripcion del pilar]",
            )
        )


def parse_topic_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if pair:
        month, value = pair
        data.monthly_topics.append((month.strip(), [item.strip() for item in value.split("|") if item.strip()]))


def parse_journalist_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    if normalize_label(label).startswith("periodista"):
        parts = [part.strip() for part in value.split("|")]
        data.journalists.append(
            Journalist(
                name=parts[0] if len(parts) > 0 else "[COMPLETAR]",
                outlet=parts[1] if len(parts) > 1 else "[COMPLETAR]",
                beat=parts[2] if len(parts) > 2 else "[COMPLETAR]",
            )
        )


def parse_step_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if pair:
        _, value = pair
        data.strategic_steps.append(value)


def parse_social_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    normalized = normalize_label(label)
    if normalized.startswith("red social"):
        parts = [part.strip() for part in value.split("|") if part.strip()]
        platform = parts[0] if parts else "[COMPLETAR: plataforma]"
        details = {
            "seguidores": "[COMPLETAR]",
            "frecuencia de posteo": "[COMPLETAR]",
            "tono": "[COMPLETAR]",
            "observaciones": "[COMPLETAR]",
        }
        for part in parts[1:]:
            if ":" not in part:
                continue
            key, current = part.split(":", 1)
            details[normalize_label(key)] = current.strip()
        data.socials.append(
            SocialChannel(
                platform=platform,
                followers=details["seguidores"],
                frequency=details["frecuencia de posteo"],
                tone=details["tono"],
                observations=details["observaciones"],
            )
        )
    elif normalized == "observaciones generales":
        data.social_summary = value


def parse_tactic_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if not pair:
        return
    label, value = pair
    if normalize_label(label).startswith("tactica"):
        parts = [part.strip() for part in value.split("|", 1)]
        data.tactics.append(
            Tactic(
                name=parts[0] if parts else "[COMPLETAR: tactica]",
                description=parts[1] if len(parts) > 1 else "[COMPLETAR: descripcion de la tactica]",
            )
        )


def parse_execution_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if pair:
        _, value = pair
        data.execution_plan = value


def parse_metric_line(line: str, data: ParsedData) -> None:
    pair = split_label_value(line)
    if pair:
        _, value = pair
        data.metrics.append(value)


def post_process(data: ParsedData) -> None:
    if not data.general.get("nombre del cliente"):
        data.pending.append("Nombre del cliente")
    if not data.general.get("objetivo central de la estrategia"):
        data.pending.append("Objetivo central de la estrategia")
    if not data.competitors:
        data.pending.append("Analisis de competidores")
    if not data.concept.get("concepto clave analizado"):
        data.pending.append("Concepto clave analizado")
    if not data.persona.archetype:
        data.pending.append("Buyer persona")
    if len(data.pillars) < 2:
        data.pending.append("Pilares de comunicacion")
    if not data.monthly_topics:
        data.pending.append("Propuestas de tematicas")
    if len(data.strategic_steps) < 5:
        data.pending.append("5 pasos estrategicos")
    if not data.socials:
        data.pending.append("Estado de redes sociales del cliente")
    if not data.tactics:
        data.pending.append("Tacticas / tipos de contenido")
    if not data.execution_plan:
        data.pending.append("Plan de ejecucion")
    if not data.metrics:
        data.pending.append("Metricas de exito")


def detect_plan_type(data: ParsedData) -> str:
    declared = normalize_label(data.general.get("tipo de plan", ""))
    if "ambos" in declared:
        return "combined"
    if "prensa" in declared and "marketing" in declared:
        return "combined"
    if "prensa" in declared:
        return "press"
    if "marketing" in declared:
        return "content"

    has_press = bool(data.concept or data.journalists or data.strategic_steps)
    has_content = bool(data.socials or data.tactics or data.execution_plan)
    if has_press and has_content:
        return "combined"
    if has_press:
        return "press"
    return "content"


def prompt_combined_mode() -> str:
    print("El documento contiene informacion para ambos planes.")
    print("1. Presentacion combinada")
    print("2. Dos presentaciones separadas")
    while True:
        answer = input("Selecciona una opcion (1 o 2): ").strip()
        if answer == "1":
            return "combined"
        if answer == "2":
            return "separate"


def slugify(value: str) -> str:
    value = normalize_label(value).replace(" ", "_")
    return re.sub(r"[^a-z0-9_]", "", value) or "cliente"


def find_latest_input() -> Path | None:
    files = [path for path in INPUT_DIR.iterdir() if path.suffix.lower() in {".txt", ".docx", ".pdf"}]
    return max(files, key=lambda current: current.stat().st_mtime) if files else None


def write_summary(data: ParsedData, pptx_path: Path, plan_type: str) -> None:
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    summary_path = pptx_path.with_suffix(".txt")
    completed = [
        "DATOS GENERALES",
        "ANALISIS DE COMPETIDORES" if data.competitors else None,
        "ANALISIS DE CONCEPTO CLAVE" if data.concept else None,
        "BUYER PERSONA" if data.persona.archetype else None,
        "PILARES DE COMUNICACION" if data.pillars else None,
        "PROPUESTAS DE TEMATICAS" if data.monthly_topics else None,
        "PERIODISTAS CLAVE" if data.journalists else None,
        "PASOS ESTRATEGICOS" if data.strategic_steps else None,
        "ESTADO DE REDES SOCIALES DEL CLIENTE" if data.socials else None,
        "TACTICAS / TIPOS DE CONTENIDO" if data.tactics else None,
        "PLAN DE EJECUCION" if data.execution_plan else None,
        "METRICAS DE EXITO" if data.metrics else None,
    ]
    lines = [f"Archivo generado: {pptx_path.name}", f"Tipo de plan: {plan_type}", "", "Secciones completadas:"]
    lines.extend(f"- {section}" for section in completed if section)
    lines.append("")
    lines.append("INFORMACION PENDIENTE DE COMPLETAR:")
    if data.pending:
        lines.extend(f"- {item}" for item in data.pending)
    else:
        lines.append("- No se detectaron campos pendientes criticos.")
    summary_path.write_text("\n".join(lines), encoding="utf-8")


def get_output_dir() -> Path:
    candidates = [
        OUTPUT_DIR,
        Path(tempfile.gettempdir()) / "strategic-ppt-generator-output",
    ]
    for candidate in candidates:
        try:
            candidate.mkdir(parents=True, exist_ok=True)
            probe = candidate / ".write_test"
            probe.write_text("ok", encoding="utf-8")
            probe.unlink(missing_ok=True)
            return candidate
        except OSError:
            continue
    fallback = Path(tempfile.mkdtemp(prefix="strategic-ppt-generator-"))
    fallback.mkdir(parents=True, exist_ok=True)
    return fallback


def add_picture_safe(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_brand_logo(slide, dark: bool = False) -> None:
    if ASSET_LOGO.exists():
        add_picture_safe(slide, ASSET_LOGO, Inches(10.85), Inches(0.2), width=Inches(1.95))
    else:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(0.2), Inches(1.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_PRIMARY
        badge.line.fill.background()
        add_text_inside_shape(badge, "PRS", COLOR_WHITE, 16, bold=True)


def add_footer_tagline(slide, dark: bool = False) -> None:
    if ASSET_TAGLINE.exists():
        add_picture_safe(slide, ASSET_TAGLINE, Inches(0.75), Inches(6.82), width=Inches(3.0))
    else:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(6.85), Inches(3.5), Inches(0.25))
        p = box.text_frame.paragraphs[0]
        p.text = "La agencia de las startups de Latam"
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE if dark else COLOR_MUTED


def add_brand_corner(slide, dark: bool = False) -> None:
    path = ASSET_BRACKET_WHITE if dark else ASSET_BRACKET_PINK
    if path.exists():
        add_picture_safe(slide, path, Inches(0.7), Inches(0.24), width=Inches(0.55))


def add_full_bleed_brand_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BRAND_DARK
    if ASSET_SKYLINE.exists():
        add_picture_safe(slide, ASSET_SKYLINE, 0, 0, width=Inches(13.333), height=Inches(7.5))
        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = COLOR_BRAND_DARK
        overlay.fill.transparency = 0.16
        overlay.line.fill.background()


def add_slide_base(prs: Presentation, title: str, dark: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        add_full_bleed_brand_background(slide)
        title_left = Inches(0.95)
        title_top = Inches(0.52)
        title_color = COLOR_WHITE
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BRAND_PALE
        top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.05))
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = COLOR_BRAND_DARK
        top_band.line.fill.background()
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.26), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLOR_PRIMARY
        accent.line.fill.background()
        title_left = Inches(0.88)
        title_top = Inches(0.28)
        title_color = COLOR_WHITE
    title_box = slide.shapes.add_textbox(title_left, title_top, Inches(9.5), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Trebuchet MS"
    run.font.bold = True
    run.font.size = Pt(29 if not dark else 34)
    run.font.color.rgb = title_color
    add_brand_corner(slide, dark=dark)
    add_brand_logo(slide, dark=dark)
    if not dark:
        add_footer_tagline(slide, dark=False)
    return slide


def add_text_inside_shape(shape, text: str, color: RGBColor, size: int, bold: bool = False) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_block(slide, title: str, body: str, left, top, width, height) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_CARD_LINE
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Trebuchet MS"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BRAND_DARK
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Calibri"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_DARK
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.14), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_PRIMARY
    accent.line.fill.background()


def add_card_heading(slide, text: str, left, top, width) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Trebuchet MS"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BRAND_DARK


def add_simple_text(slide, text: str, left, top, width, height) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_DARK


def add_big_stat(slide, value: str, left, top, label: str) -> None:
    box = slide.shapes.add_textbox(left, top, Inches(2.0), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = value
    p.font.name = "Trebuchet MS"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_BRAND_DARK
    p2 = box.text_frame.add_paragraph()
    p2.text = label
    p2.font.name = "Calibri"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_MUTED


def add_chip_group(slide, items: Iterable[str], left, top, width) -> None:
    current_x = left
    current_y = top
    max_x = left + width
    for item in items:
        chip_width = max(Inches(1.1), Inches(min(2.6, 0.5 + len(item) * 0.06)))
        if current_x + chip_width > max_x:
            current_x = left
            current_y += Inches(0.5)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, current_x, current_y, chip_width, Inches(0.35))
        chip.fill.solid()
        chip.fill.fore_color.rgb = COLOR_PRIMARY
        chip.line.fill.background()
        add_text_inside_shape(chip, item, COLOR_WHITE, 10)
        current_x += chip_width + Inches(0.12)


def add_pie_chart(slide, tiers: dict[str, int], left, top, width, height) -> None:
    chart_data = CategoryChartData()
    categories = ["Tier 1", "Tier 2", "Tier 3", "Tier 4"]
    values = [max(0, tiers.get(category, 0)) for category in categories]
    if sum(values) == 0:
        values = [1, 1, 1, 1]
    chart_data.categories = categories
    chart_data.add_series("Distribucion", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, chart_data).chart
    chart.has_legend = True
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series = chart.series[0]
    for index, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = TIER_COLORS[categories[index]]
        point.format.line.color.rgb = COLOR_WHITE


def add_cover_slide(prs: Presentation, title: str, data: ParsedData) -> None:
    slide = add_slide_base(prs, title, dark=True)
    client = data.general.get("nombre del cliente", "[COMPLETAR: nombre del cliente]")
    agency = data.general.get("agencia", "")

    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.85), Inches(6.15), Inches(2.45))
    block.fill.solid()
    block.fill.fore_color.rgb = COLOR_WHITE
    block.fill.transparency = 0.08
    block.line.fill.background()

    subtitle = slide.shapes.add_textbox(Inches(1.1), Inches(2.65), Inches(6.0), Inches(0.8))
    p = subtitle.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{client} | {datetime.now().year}"
    run.font.name = "Trebuchet MS"
    run.font.bold = True
    run.font.size = Pt(21)
    run.font.color.rgb = COLOR_WHITE

    if agency:
        footer = slide.shapes.add_textbox(Inches(1.1), Inches(5.9), Inches(10.5), Inches(0.4))
        p = footer.text_frame.paragraphs[0]
        p.text = agency
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_WHITE

    kicker = slide.shapes.add_textbox(Inches(1.1), Inches(2.18), Inches(4.5), Inches(0.35))
    p = kicker.text_frame.paragraphs[0]
    p.text = "REPORTE ESTRATEGICO"
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_WHITE
    add_footer_tagline(slide, dark=True)


def add_section_divider(prs: Presentation, title: str) -> None:
    slide = add_slide_base(prs, title, dark=True)
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.2), Inches(5.35), Inches(1.45))
    block.fill.solid()
    block.fill.fore_color.rgb = COLOR_WHITE
    block.fill.transparency = 0.08
    block.line.fill.background()
    add_text_inside_shape(block, title, COLOR_WHITE, 24, bold=True)
    add_footer_tagline(slide, dark=True)


def add_close_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Gracias", dark=True)
    box = slide.shapes.add_textbox(Inches(0.95), Inches(2.35), Inches(9.5), Inches(1.3))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = data.general.get("agencia", "Presentacion lista para revision")
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = COLOR_WHITE
    p2 = box.text_frame.add_paragraph()
    p2.text = "Conoce mas de nosotros"
    p2.font.name = "Trebuchet MS"
    p2.font.bold = True
    p2.font.size = Pt(30)
    p2.font.color.rgb = COLOR_WHITE
    add_footer_tagline(slide, dark=True)


def add_competitor_slides(prs: Presentation, data: ParsedData, title: str = "Analisis de Competidores") -> None:
    competitors = data.competitors or [Competitor("[COMPLETAR: competidor]", "[COMPLETAR]", {})]
    for start in range(0, len(competitors), 3):
        slide = add_slide_base(prs, title)
        chunk = competitors[start : start + 3]
        for index, competitor in enumerate(chunk):
            left = Inches(0.95 + index * 3.95)
            top = Inches(1.35)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(3.55), Inches(4.35))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_WHITE
            card.line.color.rgb = COLOR_CARD_LINE
            add_card_heading(slide, competitor.name, left + Inches(0.15), top + Inches(0.1), Inches(3.0))
            add_big_stat(slide, competitor.total, left + Inches(0.15), top + Inches(0.55), "Total de notas")
            add_pie_chart(slide, competitor.tiers, left + Inches(0.15), top + Inches(1.55), Inches(3.1), Inches(2.1))

        add_text_block(
            slide,
            "Conclusiones",
            data.competitor_conclusions or "[CONCLUSIONES DEL ANALISIS DE COMPETIDORES]",
            Inches(0.95),
            Inches(5.95),
            Inches(11.3),
            Inches(0.8),
        )


def add_concept_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Analisis de Concepto Clave")
    add_card_heading(
        slide,
        data.concept.get("concepto clave analizado", "[COMPLETAR: concepto clave]"),
        Inches(0.95),
        Inches(1.2),
        Inches(4.6),
    )
    add_big_stat(
        slide,
        data.concept.get("total de notas del concepto", "[COMPLETAR]"),
        Inches(0.95),
        Inches(1.7),
        "Total de notas",
    )
    add_pie_chart(
        slide,
        {
            "Tier 1": safe_int(data.concept.get("tier 1", "0")),
            "Tier 2": safe_int(data.concept.get("tier 2", "0")),
            "Tier 3": safe_int(data.concept.get("tier 3", "0")),
            "Tier 4": safe_int(data.concept.get("tier 4", "0")),
        },
        Inches(0.95),
        Inches(2.55),
        Inches(3.7),
        Inches(2.5),
    )
    add_chip_group(
        slide,
        data.concept_themes or ["[COMPLETAR: tematicas mas repetidas]"],
        Inches(5.0),
        Inches(1.7),
        Inches(6.5),
    )
    add_text_block(
        slide,
        "Conclusiones",
        data.concept.get("conclusiones", "[CONCLUSIONES DEL CONCEPTO CLAVE]"),
        Inches(5.0),
        Inches(4.55),
        Inches(6.5),
        Inches(1.4),
    )


def add_conclusions_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Conclusiones Generales")
    insights = split_multi_value(data.competitor_conclusions or "")[:3]
    while len(insights) < 3:
        insights.append("[COMPLETAR: conclusiones estrategicas]")
    for index, insight in enumerate(insights):
        left = Inches(0.95 + index * 3.95)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.45), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_CARD_LINE
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, left + Inches(0.18), Inches(2.05), Inches(0.55), Inches(0.55))
        icon.fill.solid()
        icon.fill.fore_color.rgb = COLOR_PRIMARY
        icon.line.fill.background()
        add_simple_text(slide, insight, left + Inches(0.2), Inches(2.8), Inches(2.95), Inches(1.7))


def add_resources_slide(prs: Presentation) -> None:
    slide = add_slide_base(prs, "Recursos de Prensa")
    resources = [
        ("Comunicados", "Comunicados, columnas de opinion o cartas al director. Es el recurso mas usado para lograr apariciones en medios digitales."),
        ("Entrevistas", "Centradas en potenciar estudios o hitos. Acompanadas de un brief enfocado en la estrategia."),
        ("Contingencia", "Participacion en solicitudes de periodistas. Contenido propuesto segun pauta mediatica del momento."),
        ("Relacionamiento", "Instancias de conversacion intima con medios para crear lazos a largo plazo con la marca."),
    ]
    for index, (title, desc) in enumerate(resources):
        left = Inches(0.95 + (index % 2) * 5.85)
        top = Inches(1.6 + (index // 2) * 2.1)
        add_text_block(slide, title, desc, left, top, Inches(5.25), Inches(1.65))


def add_objective_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Objetivo Core de la Estrategia")
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(1.85), Inches(10.95), Inches(3.2))
    block.fill.solid()
    block.fill.fore_color.rgb = COLOR_PRIMARY
    block.line.fill.background()
    tf = block.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.general.get("objetivo central de la estrategia", "[COMPLETAR: objetivo central de la estrategia]")
    run.font.name = "Trebuchet MS"
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = COLOR_WHITE


def add_buyer_persona_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Buyer Persona")
    persona = data.persona
    add_card_heading(slide, persona.archetype or "[COMPLETAR: nombre o arquetipo]", Inches(0.95), Inches(1.15), Inches(4.5))
    blocks = [
        ("Edad / Profesion", f"{persona.age or '[COMPLETAR]'} | {persona.profession or '[COMPLETAR]'}", 0.95, 1.7, 3.7, 1.0),
        ("Contexto", persona.context or "[COMPLETAR: contexto]", 0.95, 2.95, 5.65, 1.15),
        ("Dolores", persona.pains or "[COMPLETAR: dolores]", 6.8, 1.7, 5.0, 1.15),
        ("Metas", persona.goals or "[COMPLETAR: metas y objetivos]", 0.95, 4.3, 5.65, 1.15),
        ("Como busca informacion", persona.info_behavior or "[COMPLETAR: comportamiento para buscar informacion]", 6.8, 3.15, 5.0, 2.3),
    ]
    for title, body, left, top, width, height in blocks:
        add_text_block(slide, title, body, Inches(left), Inches(top), Inches(width), Inches(height))
    add_chip_group(slide, persona.personality or ["[COMPLETAR: personalidad]"], Inches(6.8), Inches(1.95), Inches(5.0))


def add_steps_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "5 Pasos de la Estrategia")
    steps = list(data.strategic_steps[:5])
    while len(steps) < 5:
        steps.append(f"[COMPLETAR: Paso {len(steps) + 1}]")
    for index, step in enumerate(steps):
        top = Inches(1.35 + index * 1.05)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.0), top, Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_PRIMARY
        badge.line.fill.background()
        add_text_inside_shape(badge, str(index + 1), COLOR_WHITE, 14, bold=True)
        add_simple_text(slide, step, Inches(1.8), top - Inches(0.02), Inches(9.9), Inches(0.65))


def add_pillars_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Pilares de Comunicacion")
    pillars = list(data.pillars[:2])
    while len(pillars) < 2:
        pillars.append(Pillar("[COMPLETAR: nombre del pilar]", "[COMPLETAR: descripcion del pilar]"))
    for index, pillar in enumerate(pillars):
        left = Inches(0.95 + index * 5.95)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.85), Inches(5.35), Inches(3.45))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BRAND_DARK if index == 0 else COLOR_WHITE
        card.line.color.rgb = COLOR_CARD_LINE
        tf = card.text_frame
        p = tf.paragraphs[0]
        p.text = pillar.name
        p.font.name = "Trebuchet MS"
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_WHITE if index == 0 else COLOR_BRAND_DARK
        p2 = tf.add_paragraph()
        p2.text = pillar.description
        p2.font.name = "Calibri"
        p2.font.size = Pt(15)
        p2.font.color.rgb = COLOR_WHITE if index == 0 else COLOR_DARK


def add_topics_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Propuestas de Tematicas")
    topics = data.monthly_topics or [("Mes 1", ["[COMPLETAR: titular 1]", "[COMPLETAR: titular 2]"])]
    for index, (month, items) in enumerate(topics[:6]):
        left = Inches(0.95 + (index % 2) * 5.95)
        top = Inches(1.45 + (index // 2) * 1.75)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_PRIMARY if index % 2 == 0 else COLOR_WHITE
        card.line.color.rgb = COLOR_CARD_LINE
        tf = card.text_frame
        p = tf.paragraphs[0]
        p.text = month
        p.font.name = "Trebuchet MS"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_WHITE if index % 2 == 0 else COLOR_BRAND_DARK
        for item in (items[:2] or ["[COMPLETAR: titular]"]):
            p2 = tf.add_paragraph()
            p2.text = f"- {item}"
            p2.font.name = "Calibri"
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLOR_WHITE if index % 2 == 0 else COLOR_DARK


def add_journalists_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Periodistas Clave")
    headers = ["Nombre", "Medio", "Area / Beat"]
    lefts = [0.95, 4.2, 7.55]
    widths = [3.0, 3.1, 3.7]
    for left, width, header in zip(lefts, widths, headers):
        header_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(1.35), Inches(width), Inches(0.45))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_BRAND_DARK
        header_shape.line.fill.background()
        add_text_inside_shape(header_shape, header, COLOR_WHITE, 13, bold=True)
    journalists = list(data.journalists[:10])
    while len(journalists) < 10:
        journalists.append(Journalist("[COMPLETAR]", "[COMPLETAR]", "[COMPLETAR]"))
    for index, journalist in enumerate(journalists):
        top = Inches(1.85 + index * 0.45)
        values = [journalist.name, journalist.outlet, journalist.beat]
        for left, width, value in zip(lefts, widths, values):
            row = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), top, Inches(width), Inches(0.42))
            row.fill.solid()
            row.fill.fore_color.rgb = COLOR_WHITE if index % 2 == 0 else COLOR_BRAND_PALE
            row.line.color.rgb = COLOR_CARD_LINE
            add_text_inside_shape(row, value, COLOR_DARK, 11)


def add_social_status_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Desk Research: Estado de Redes Sociales del Cliente")
    socials = data.socials or [SocialChannel("[COMPLETAR]", "[COMPLETAR]", "[COMPLETAR]", "[COMPLETAR]", "[COMPLETAR]")]
    for index, social in enumerate(socials[:4]):
        left = Inches(0.95 + (index % 2) * 5.9)
        top = Inches(1.45 + (index // 2) * 2.05)
        body = (
            f"Seguidores: {social.followers}\n"
            f"Frecuencia: {social.frequency}\n"
            f"Tono: {social.tone}\n"
            f"Observaciones: {social.observations}"
        )
        add_text_block(slide, social.platform, body, left, top, Inches(5.3), Inches(1.65))
    add_text_block(
        slide,
        "Sintesis general",
        data.social_summary or "[COMPLETAR: observaciones generales de redes sociales]",
        Inches(0.95),
        Inches(5.9),
        Inches(11.2),
        Inches(0.7),
    )


def add_goals_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Objetivos / Metas")
    add_text_block(
        slide,
        "Objetivo principal",
        data.general.get("objetivo central de la estrategia", "[COMPLETAR: objetivo central]"),
        Inches(0.95),
        Inches(1.4),
        Inches(11.2),
        Inches(1.0),
    )
    metrics = data.metrics[:4] if data.metrics else ["[COMPLETAR: objetivo o KPI]"]
    for index, item in enumerate(metrics):
        left = Inches(0.95 + (index % 2) * 5.9)
        top = Inches(2.8 + (index // 2) * 1.2)
        add_text_block(slide, f"Meta {index + 1}", item, left, top, Inches(5.3), Inches(0.9))


def add_tactics_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Tacticas / Tipos de Contenido")
    tactics = data.tactics or [Tactic("[COMPLETAR: tactica]", "[COMPLETAR: descripcion]")]
    for index, tactic in enumerate(tactics[:6]):
        left = Inches(0.95 + (index % 2) * 5.9)
        top = Inches(1.45 + (index // 2) * 1.25)
        add_text_block(slide, tactic.name, tactic.description, left, top, Inches(5.3), Inches(0.95))


def add_execution_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Como Se Ejecutara")
    add_text_block(
        slide,
        "Roadmap de ejecucion",
        data.execution_plan or "[COMPLETAR: descripcion de como se ejecutara la estrategia]",
        Inches(0.95),
        Inches(1.5),
        Inches(11.2),
        Inches(2.0),
    )
    for index, stage in enumerate(["Planificacion", "Produccion", "Distribucion", "Medicion"]):
        left = Inches(0.95 + index * 2.8)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(4.35), Inches(2.3), Inches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_PRIMARY if index % 2 == 0 else COLOR_WHITE
        card.line.color.rgb = COLOR_CARD_LINE
        add_text_inside_shape(card, stage, COLOR_WHITE if index % 2 == 0 else COLOR_BRAND_DARK, 15, bold=True)


def add_metrics_slide(prs: Presentation, data: ParsedData) -> None:
    slide = add_slide_base(prs, "Como Mediremos el Exito?")
    metrics = data.metrics or ["[COMPLETAR: KPIs de exito]"]
    for index, metric in enumerate(metrics[:6]):
        left = Inches(0.95 + (index % 2) * 5.9)
        top = Inches(1.45 + (index // 2) * 1.2)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_CARD_LINE
        add_text_inside_shape(card, metric, COLOR_DARK, 13)


def build_press_deck(prs: Presentation, data: ParsedData, include_cover: bool) -> None:
    if include_cover:
        add_cover_slide(prs, "Plan Estrategico de Prensa", data)
    add_competitor_slides(prs, data)
    add_concept_slide(prs, data)
    add_conclusions_slide(prs, data)
    add_resources_slide(prs)
    add_objective_slide(prs, data)
    add_buyer_persona_slide(prs, data)
    add_steps_slide(prs, data)
    add_pillars_slide(prs, data)
    add_topics_slide(prs, data)
    add_journalists_slide(prs, data)
    add_metrics_slide(prs, data)


def build_content_deck(prs: Presentation, data: ParsedData, include_cover: bool) -> None:
    if include_cover:
        add_cover_slide(prs, "Plan de Marketing de Contenidos", data)
    add_competitor_slides(prs, data, title="Analisis de Competidores")
    add_social_status_slide(prs, data)
    add_goals_slide(prs, data)
    add_tactics_slide(prs, data)
    add_execution_slide(prs, data)
    add_metrics_slide(prs, data)


def build_combined_deck(prs: Presentation, data: ParsedData) -> None:
    add_cover_slide(prs, "Plan Estrategico de Comunicaciones", data)
    add_section_divider(prs, "Plan de Prensa")
    build_press_deck(prs, data, include_cover=False)
    add_section_divider(prs, "Plan de Marketing de Contenidos")
    build_content_deck(prs, data, include_cover=False)
    add_close_slide(prs, data)


def save_presentation(data: ParsedData, plan_type: str, source_path: Path, output_dir: Path | None = None) -> Path:
    output_dir = output_dir or get_output_dir()
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if plan_type == "press":
        build_press_deck(prs, data, include_cover=True)
        prefix = "Plan_Estrategico"
    elif plan_type == "content":
        build_content_deck(prs, data, include_cover=True)
        prefix = "Plan_Marketing_Contenidos"
    else:
        build_combined_deck(prs, data)
        prefix = "Plan_Estrategico"

    client = slugify(data.general.get("nombre del cliente", source_path.stem))
    output_path = output_dir / f"{prefix}_{client}_{datetime.now().year}.pptx"
    prs.save(output_path)
    write_summary(data, output_path, plan_type)
    return output_path


def build_presentation(data: ParsedData, plan_type: str, source_path: Path, output_dir: Path | None = None) -> list[Path]:
    output_dir = output_dir or get_output_dir()
    if plan_type == "separate":
        return [
            save_presentation(data, "press", source_path, output_dir=output_dir),
            save_presentation(data, "content", source_path, output_dir=output_dir),
        ]
    return [save_presentation(data, plan_type, source_path, output_dir=output_dir)]


def generate_from_file(source: Path, mode: str | None = None, prompt_on_combined: bool = True, output_dir: Path | None = None) -> tuple[list[Path], ParsedData, str]:
    source = source.expanduser().resolve()
    data = parse_document(extract_text(source))
    plan_type = mode or detect_plan_type(data)
    if plan_type == "combined" and not mode and prompt_on_combined:
        plan_type = prompt_combined_mode()
    outputs = build_presentation(data, plan_type, source, output_dir=output_dir)
    return outputs, data, plan_type


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Genera presentaciones estrategicas en PPTX.")
    parser.add_argument("input_file", nargs="?", help="Ruta del documento de entrada")
    parser.add_argument("--auto", action="store_true", help="Usa automaticamente el archivo mas reciente en input")
    parser.add_argument("--mode", choices=["press", "content", "combined", "separate"], help="Fuerza un modo de salida")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    OUTPUT_DIR.mkdir(exist_ok=True)

    if args.input_file:
        source = Path(args.input_file).expanduser().resolve()
    elif args.auto:
        source = find_latest_input()
        if source is None:
            print("No se encontro ningun archivo compatible dentro de la carpeta input.")
            return 1
    else:
        print("Debes indicar un archivo de entrada o usar --auto.")
        return 1

    if not source.exists():
        print(f"No se encontro el archivo: {source}")
        return 1

    outputs, data, plan_type = generate_from_file(source, mode=args.mode, prompt_on_combined=not args.mode)
    print("")
    print("Archivos generados:")
    for output in outputs:
        print(f"- {output}")
        print(f"- {output.with_suffix('.txt')}")
    if data.pending:
        print("")
        print("INFORMACION PENDIENTE DE COMPLETAR:")
        for item in data.pending:
            print(f"- {item}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
